
from django.shortcuts import render,redirect,get_object_or_404
from .models import UploadedFile
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.contrib.auth import logout as django_logout
import os
import io
import fitz  # PyMuPDF for PDF parsing
import docx  # python-docx for .docx parsing
# import textract  
import csv
from pptx import Presentation
from PIL import Image
from django.http import HttpResponseBadRequest
from django.http import HttpResponse
from .models import ChatSession,Feedback,SimilarQuestion
from sklearn.metrics.pairwise import cosine_similarity
from django.urls import reverse_lazy
import numpy as np
import nltk
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
import inflect
nltk.download('punkt')
nltk.download('wordnet')
import re
from sentence_transformers import SentenceTransformer
from .forms import QnAForm
from django.contrib.auth.decorators import login_required
from django.contrib.admin.views.decorators import staff_member_required
import nltk
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from nltk.stem import WordNetLemmatizer
from django.utils.crypto import get_random_string
from bs4 import BeautifulSoup
import requests
from django.contrib.auth.forms import AuthenticationForm
from django.contrib.auth import authenticate, login
from django.utils import timezone
import uuid 

import fitz
from .decorators import api_key_required


nltk.download('punkt')
nltk.download('stopwords')
nltk.download('wordnet')



from django.urls import reverse




def index(request):
    return render(request,'index.html')


 

@csrf_exempt
@api_key_required
@login_required(login_url='/admin_login/')
def upload_file(request):
    if request.method == 'POST':
        if 'file' in request.FILES:             
            uploaded_file = request.FILES['file']
            current_user = request.user
            chatbot_name = request.POST.get('bot_name')
            chatbot_description = request.POST.get('bot_desc')

            # Save the uploaded file to the database
            uploaded_file_obj = UploadedFile.objects.create(file=uploaded_file, user=current_user,chatbot_name=chatbot_name,chatbot_description=chatbot_description)

            # Get the ID of the uploaded file
            file_id = uploaded_file_obj.id

            # Redirect to the chat page after uploading the file, including the file ID in the URL
            return redirect(reverse('chat', kwargs={'user_name': current_user, 'file_id': file_id})+ '?success=true')
        elif 'link' in request.POST:
            link = request.POST['link']
            print(link)
            current_user = request.user
            chatbot_name = request.POST.get('bot_name')
            chatbot_description = request.POST.get('bot_desc')

            
            response = requests.get(link)
            if response.status_code == 200:
                # If the request is successful, parse the content
                soup = BeautifulSoup(response.content, 'html.parser')
                paragraphs = soup.find_all('p')
                content = '\n'.join([p.get_text() for p in paragraphs])
                
                # Save the link and content to the database
                uploaded_file_obj = UploadedFile.objects.create(link=link, content=content, user=current_user,chatbot_name=chatbot_name,chatbot_description=chatbot_description)
                file_id = uploaded_file_obj.id
            

            return redirect(reverse('chat', kwargs={'user_name': current_user, 'file_id': file_id})+ '?success=true')
            # else:
            #     # If the request fails, handle the error
            #     return HttpResponseBadRequest("Failed to fetch content from the provided URL.")

        else:
            # If neither a file nor a link is provided, return a bad request response
            return HttpResponseBadRequest("Invalid request. Please provide either a file or a link.")
    
    return render(request, 'upload_file.html')


@api_key_required
def custom_login_view(request):
    # Redirect authenticated users
    if request.user.is_authenticated:
        if request.user.is_superuser:
            return HttpResponseRedirect(reverse_lazy("dashboard"))
        else:
            return HttpResponseRedirect(reverse_lazy("unauthorised"))

    if request.method == "POST":
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            user = authenticate(username=form.cleaned_data['username'], password=form.cleaned_data['password'])
            if user is not None:
                login(request, user)
                return redirect(get_success_url(request))
            else:
                # Handle the case where authentication fails
                return render(request, "admin_login.html", {'form': form})
        else:
            # Form is invalid
            return render(request, "admin_login.html", {'form': form})
    else:
        # GET request, show the empty form
        form = AuthenticationForm()
        return render(request, "admin_login.html", {'form': form})


@api_key_required
def get_success_url(request):
    if request.user.is_authenticated and request.user.is_superuser:
        return reverse_lazy("dashboard")
    else:
        return reverse_lazy("unauthorised")


@api_key_required
@login_required(login_url='/admin_login/')
def dashboard_view(request):
    if not request.user.is_superuser:
        return redirect("unauthorised")
    
    uploaded_files = UploadedFile.objects.filter(user=request.user)
    for uploaded_file in uploaded_files:
        uploaded_file.chat_url = request.build_absolute_uri(
            reverse('chat', kwargs={'user_name': request.user.username, 'file_id': uploaded_file.id})
        )
        # uploaded_file.chat_url = reverse('chat', kwargs={'user_name': request.user.username, 'file_id': uploaded_file.id})
        print(uploaded_file.chat_url)  # Debug print
    
    return render(request, 'dashboard.html', {'uploaded_files': uploaded_files})


@api_key_required
@login_required(login_url='/admin_login/')
def chatbot_details(request, file_id):
    if not request.user.is_superuser:
        return redirect("unauthorised")
    
    uploaded_file = get_object_or_404(UploadedFile, user=request.user, id=file_id)
    
    return render(request, 'details.html', {'uploaded_file': uploaded_file})



@api_key_required
@login_required(login_url='/admin_login/')
def unique_questions_table(request):
    if not request.user.is_superuser:
        return redirect("unauthorised")
        
    
    user_entries = ChatSession.objects.filter(user=request.user).order_by('-session_start_time')
    
    session_start_time = timezone.now().strftime("%B %d, %Y, %I:%M %p")
    answers_by_user = {}
    global_displayed_files = set()  # Global set to track displayed files
    
    for entry in user_entries:
        user_id = entry.user_identifier
        if user_id not in answers_by_user:
            answers_by_user[user_id] = {'answers': [], 'displayed_files': set()}
        
        if entry.file_id not in global_displayed_files:  
            global_displayed_files.add(entry.file_id)  
            try:
                uploaded_file = UploadedFile.objects.get(pk=entry.file_id)
                entry.file_name = uploaded_file.file.name
                entry.chatbot_name = UploadedFile.chatbot_name
                answers_by_user[user_id]['answers'].append(entry)
            except UploadedFile.DoesNotExist:
                continue  

    context = {
        'answers_by_user': answers_by_user,
        'session_start_time': session_start_time,
    }
    
    return render(request, 'unique_questions.html', context)



@api_key_required
@login_required(login_url='/admin_login/')
def file_queries(request, file_id):
    # uploaded_file = get_object_or_404(UploadedFile, pk=file_id)
    user_entries = ChatSession.objects.filter(user=request.user,file_id=file_id).order_by('-session_start_time')
    
    session_start_time = timezone.now().strftime("%B %d, %Y, %I:%M %p")
    answers_by_user = {}
    
    for entry in user_entries:
        user_id = entry.user_identifier
        
        if user_id not in answers_by_user:
            answers_by_user[user_id] = {'answers': []}
        
        answers_by_user[user_id]['answers'].append(entry)


        uploaded_file = get_object_or_404(UploadedFile, pk=entry.file_id)
        
        # Append the file name to the entry
        entry.file_name = uploaded_file.file.name

    
    
    context = {
        'answers_by_user': answers_by_user,
        'session_start_time': session_start_time
    }
    return render(request, 'file_details.html', context)

from django.utils import timezone

@api_key_required
@login_required(login_url='/admin_login/')
def file_queries_log(request, file_id):
    # uploaded_file = get_object_or_404(UploadedFile, pk=file_id)
    user_entries = SimilarQuestion.objects.filter(user_name=request.user,file_id=file_id).order_by('-session_start_time')
    
    session_start_time = timezone.now().strftime("%B %d, %Y, %I:%M %p")
    answers_by_user = {}
    
    for entry in user_entries:
        user_id = entry.user_session_id
        
        if user_id not in answers_by_user:
            answers_by_user[user_id] = {'answers': []}
        
        answers_by_user[user_id]['answers'].append(entry)


        uploaded_file = get_object_or_404(UploadedFile, pk=entry.file_id)
        
        # Append the file name to the entry
        entry.file_name = uploaded_file.file.name

    
    
    context = {
        'answers_by_user': answers_by_user,
        'session_start_time': session_start_time
    }
    return render(request, 'file_details_log.html', context)








@api_key_required
@login_required(login_url='/admin_login/')
def dashboard_log_question(request):
    if not request.user.is_superuser:
        return redirect("unauthorised")
    
    user_entries = SimilarQuestion.objects.filter(user_name=request.user).order_by('-session_start_time')
    
    session_start_time = timezone.now().strftime("%B %d, %Y, %I:%M %p")
    answers_by_user = {}
    global_displayed_files = set()  # Global set to track displayed files
    
    for entry in user_entries:
        user_id = entry.user_session_id
        if user_id not in answers_by_user:
            answers_by_user[user_id] = {'answers': [], 'displayed_files': set()}
        
        if entry.file_id not in global_displayed_files:  # Check against the global set
            global_displayed_files.add(entry.file_id)  # Add to the global set
            try:
                uploaded_file = UploadedFile.objects.get(pk=entry.file_id)
                entry.file_name = uploaded_file.file.name
                answers_by_user[user_id]['answers'].append(entry)
            except UploadedFile.DoesNotExist:
                continue  # Skip adding entry if the file does not exist

    context = {
        'answers_by_user': answers_by_user,
        'session_start_time': session_start_time,
    }
    
    return render(request, 'log_questions.html', context)


@api_key_required
def logout(request):
    django_logout(request)
    return redirect('admin_login')



@api_key_required
def qna_create(request):
    if request.method == 'POST':
        form = QnAForm(request.POST)
        if form.is_valid():
            form.save()
            return redirect('dashboard')
    else:
        form = QnAForm()  # Initialize the form for GET requests
    return render(request, 'qna_form.html', {'form': form})


@api_key_required
def qna_update(request, pk):
    qna = get_object_or_404(ChatSession, pk=pk)
    if request.method == 'POST':
        form = QnAForm(request.POST, instance=qna)
        if form.is_valid():
            form.save()
            return redirect('file_details', file_id=qna.file_id)

    else:
        form = QnAForm(instance=qna)
    return render(request, 'qna_form.html', {'form': form, 'qna': qna}) 




@api_key_required
def qna_delete(request, pk):
    qna = get_object_or_404(ChatSession, pk=pk)
    if request.method == 'POST':  # Corrected 'methos' to 'method'
        qna.delete()
        return redirect('file_details', file_id=qna.file_id)
    return render(request, 'qna_confirm_delete.html', {'qna': qna})



from django.contrib.auth.models import User
@api_key_required
def qna_delete_log(request, pk):
    qna = get_object_or_404(SimilarQuestion, pk=pk)
    if request.method == 'POST':  # Corrected 'methos' to 'method'
        qna.delete()
        return redirect('file_details_log_questions', file_id=qna.file_id)
    return render(request, 'qna_confirm_delete.html', {'qna': qna})


@api_key_required
def chatbot_delete(request, pk):
    qna = get_object_or_404(UploadedFile, pk=pk)
    if request.method == 'POST':  # Corrected 'methos' to 'method'
        qna.delete()
        return redirect('dashboard')
    return render(request, 'chatbot.html', {'qna': qna})


    

from django.contrib import messages
@api_key_required
@login_required(login_url='/admin_login/')
def verify_question(request, session_id, file_id):
    sessions = ChatSession.objects.filter(file_id=file_id, session_id=session_id)
    if not sessions:
        messages.error(request, 'No matching session found')
        return redirect('file_details', file_id=file_id)  # Redirect to the appropriate view

    verified = False
    for session in sessions:
        if request.user not in session.verified_by.all():
            session.verification_count += 1
            session.verified_by.add(request.user)
            session.save()
            verified = True

    if verified:
        messages.success(request, 'Verification successful')
    else:
        messages.info(request, 'Already verified by this user')

    return redirect('file_details', file_id=file_id)








import os
import fitz  # PyMuPDF
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import numpy as np
import json
import boto3





MAX_TOKENS_USAGE = 9000
# MAX_QUERIES_ALLOWED = 1
import PyPDF2
from django.conf import settings
import tempfile

@api_key_required
def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        doc = fitz.open(pdf_path)
        for page in doc:
            text += page.get_text()
    except Exception as e:
        print(f"Error occurred while extracting text from PDF: {e}")
        raise e
    return text
@api_key_required
def split_text(text, max_length=8000):
    sentences = text.split('. ')
    segments = []
    current_segment = ""
    
    for sentence in sentences:
        if len(current_segment) + len(sentence) <= max_length:
            current_segment += sentence + '. '
        else:
            segments.append(current_segment.strip())
            current_segment = sentence + '. '
    
    if current_segment:
        segments.append(current_segment.strip())
    
    return segments
@api_key_required
def find_relevant_segment(question, segments):
    vectorizer = TfidfVectorizer().fit_transform([question] + segments)
    vectors = vectorizer.toarray()
    
    question_vector = vectors[0]
    segment_vectors = vectors[1:]
    
    similarities = cosine_similarity([question_vector], segment_vectors)
    most_relevant_index = np.argmax(similarities)
    
    return segments[most_relevant_index]


import json
import random

# GREETINGS_INPUT = ["hi", "hello", "hey", "greetings", "good day", "good morning", "good afternoon", "good evening"]
# # Responses the system can give to greetings
# GREETINGS_RESPONSE = ["Hello!", "Hi there!", "Hello! How can I assist you today?","Good Morning", "Good Afternoon", "Good Evening"]

GREETINGS_MAP = {
    "hi": "Hi there!",
    "hello": "Hello! How can I assist you today?",
    "hey": "Hey! What can I do for you?",
    "greetings": "Greetings! How may I help?",
    "good day": "Good day to you!",
    "good morning": "Good morning! Hope you're having a great day!",
    "good afternoon": "Good afternoon! How can I assist?",
    "good evening": "Good evening! What do you need help with?",
    "how are you": "I'm just a bot, but thank you for asking! How can I assist you today?"
}
@api_key_required
@csrf_exempt
def chat(request, user_name, file_id):
    print(user_name)
    file_exists = False

    if request.method == 'POST':
        # Handle POST request
        pass
    elif request.method == 'GET':
        file_exists = UploadedFile.objects.filter(id=file_id, user__username=user_name).exists()

    if request.method == 'GET' and file_exists:
        return render(request, 'chat.html', {
            'user_name': user_name,
            'file_id': file_id,
            'file_exists': file_exists,
        })
    elif request.method == 'GET' and not file_exists:
        return HttpResponse("Link Disabled.", status=404)
    elif request.method == 'POST':
        if 'user_identifier' not in request.session:
            user_identifier = f"ch-{uuid.uuid4().hex[:2]}"
            request.session['user_identifier'] = user_identifier
        else:
            user_identifier = request.session['user_identifier']

        session_start_time = request.session.get('session_start_time')
        if not session_start_time:
            session_start_time = timezone.localtime(timezone.now())
            request.session['session_start_time'] = session_start_time.strftime("%Y-%m-%d %H:%M:%S")

        session_id = f"{user_identifier}.{uuid.uuid4().hex[:2]}"

        question = request.POST.get('question', '')
        if not question:
            return HttpResponseBadRequest('Missing or empty "question" parameter in the request.')
        
        
        # if question in GREETINGS_INPUT:
        #     random_greeting = random.choice(GREETINGS_RESPONSE)  # Select a random greeting response
        #     question_asked_time = timezone.localtime(timezone.now())
        #     return JsonResponse({
        #         'answer': random_greeting,
        #         'source': 'system',
        #         'session_id': session_id,
        #         'question_asked_time': question_asked_time.strftime("%d-%m-%Y %H:%M:%S")
                
                
        #     })

        question_asked_time = timezone.localtime(timezone.now())
        request.session['current_time'] = question_asked_time.strftime("%Y-%m-%d %H:%M:%S")
        



        if question in GREETINGS_MAP:
            response_message = GREETINGS_MAP[question]
            return JsonResponse({
                'answer': response_message,
                'source': 'system',
                'question_asked_time': question_asked_time.strftime("%d-%m-%Y %H:%M:%S")
            })
        


        

        existing_answer = get_answer_from_database(question, user_identifier, session_id, user_name, file_id, question_asked_time)

        if existing_answer:
            answer, verification_count, positive_count, negative_feedback_count, file_id, question_asked_time = existing_answer
            if verification_count is not None:
                return JsonResponse({
                    'answer': answer,
                    'source': 'database',
                    'session_id': session_id,
                    'verification_count': verification_count,
                    'positive_count': positive_count,
                    'negative_feedback_count': negative_feedback_count,
                    'question_asked_time': question_asked_time.strftime("%d-%m-%Y %H:%M:%S")
                })
            else:
                return JsonResponse({
                    'answer': answer,
                    'source': 'database',
                    'session_id': session_id
                })

        user = User.objects.get(username=user_name)
        uploaded_file = get_object_or_404(UploadedFile, id=file_id, user__username=user_name)

        if uploaded_file and uploaded_file.link:
            link = uploaded_file.link
            try:
                response = requests.get(link)
                if response.status_code == 200:
                    if link.endswith(".pdf"):
                        filename = link.split("/")[-1]
                        pdf_path = os.path.join(os.getcwd(), filename)

                        with open(pdf_path, 'wb') as pdf_file:
                            pdf_file.write(response.content)

                        document_text = extract_text_from_pdf(pdf_path)
                        os.remove(pdf_path)
                    else:
                        soup = BeautifulSoup(response.content, 'html.parser')
                        paragraphs = soup.find_all('p')
                        document_text = '\n'.join([p.get_text() for p in paragraphs])
                else:
                    return JsonResponse({'error': 'Failed to fetch content from the provided link.'})
            except Exception as e:
                return JsonResponse({'error': str(e)})
        elif uploaded_file and uploaded_file.file:
            file_path = uploaded_file.file.path
            file_extension = os.path.splitext(uploaded_file.file.name)[1].lower()

            if file_extension == '.pdf':
                document_text = extract_text_from_pdf(file_path)
            elif file_extension == '.docx':
                doc = docx.Document(file_path)
                document_text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            elif file_extension == '.pptx':
                prs = Presentation(file_path)
                document_text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            elif file_extension == '.txt':
                with open(file_path, 'r') as file:
                    document_text = file.read()
            elif file_extension == '.csv':
                with open(file_path, newline='', encoding='utf-8') as csvfile:
                    csvreader = csv.reader(csvfile)
                    document_text = "\n".join([",".join(row) for row in csvreader])
            # elif file_extension == '.gif':
            #     image = Image.open(file_path)
            #     document_text = pytesseract.image_to_string(image)
            # else:
            #     document_text = textract.process(file_path, method='text', encoding='utf-8')
        else:
            return JsonResponse({'error': 'No uploaded file found.'})

        text_segments = split_text(document_text)

        

        from transformers import GPT2Tokenizer

        tokenizer = GPT2Tokenizer.from_pretrained('gpt2')

        def count_tokens(text):
            tokens = tokenizer.encode(text, add_special_tokens=False)
            return len(tokens)

        def clean_response(response_text, stop_sequences=["<end>", "</s>", "</n>"]):
            # Find the earliest occurrence of any stop sequence
            first_stop_index = min((response_text.find(seq) if seq in response_text else len(response_text) for seq in stop_sequences))
            
            # Truncate the response at the first stop sequence
            response_text = response_text[:first_stop_index].strip()

            # Remove any known unwanted tokens
            unwanted_tokens = ["ASSISTANT:", "assistant:", "user:"]
            for token in unwanted_tokens:
                response_text = response_text.replace(token, "").strip()

            return response_text

        from transformers import GPT2Tokenizer

        question_token_count = count_tokens(question)
        print(f"Token count for the question: {question_token_count}")

# Initialize the Bedrock client globally
        bedrock_client = boto3.client(service_name='bedrock-runtime')

        # Initialize the tokenizer
        tokenizer = GPT2Tokenizer.from_pretrained('gpt2')

       

        def initialize_bedrock_client():
            # Create a Bedrock client
            return boto3.client(service_name='bedrock-runtime')
        
        def truncate_context(context, max_tokens):
            tokens = tokenizer.encode(context, add_special_tokens=False)
            if len(tokens) > max_tokens:
                truncated_tokens = tokens[:max_tokens]
                return tokenizer.decode(truncated_tokens)
            return context

        def get_response(file_id, relevant_text,context, question, stop_sequences=["<end>", "</s>", "</n>"], max_attempts=3):
            print(relevant_text)
            max_context_tokens = 8192 - 1000  # Reserve tokens for the user input and response
            truncated_context = truncate_context(context, max_context_tokens)
            print(truncated_context)
            messages = [
                {"role": "system", "content": "You are a helpful assistant. Provide a brief, concise answer using only the uploaded document's content."},
                {"role": "system", "content": "You are a helpful assistant. Answer the following question based solely on the uploaded document content. Be brief and terminate your response with '<end>'."},
                {"role": "system", "content": "You are a helpful assistant. Always add a full stop at the end of the response"},
                {"role": "system","content": "If the question is about politicians, the film industry, cricketers, singers, or languages, and the specific name asked about is not in the uploaded document, respond with 'I don't know' and terminate your response with '<end>'. If the name is in the document, provide the relevant information."},
                {"role": "user", "content": f"Document details (File ID: {file_id}):"},
                {"role": "user", "content": truncated_context},
                
                {"role": "user", "content": question}
            ]
            full_prompt = "\n".join([f"{msg['role']}: {msg['content']}" for msg in messages]) + "\nassistant:"

            body = json.dumps({
                "prompt": full_prompt,
                "temperature": 0.1,
                "top_p": 0.9
            })

            # modelId = 'meta.llama3-8b-instruct-v1:0'
            modelId='meta.llama3-70b-instruct-v1:0'
            accept = 'application/json'
            contentType = 'application/json'

            # Invoke the model with the input body
            response = bedrock_client.invoke_model(body=body, modelId=modelId, accept=accept, contentType=contentType)

            # Read the StreamingBody object as a string
            response_body_str = response['body'].read().decode('utf-8')
            try:
                response_text = ""
                response_body = json.loads(response_body_str)
                # print("Debug: Response Body:", response_body)  # Debugging line
                
                # Extract and return the generation text if available
                generation = response_body.get('generation')
                print(f"Generation Text :{generation}")
                response_text += generation
                clean_resp = clean_response(response_text, stop_sequences)
                token_count = count_tokens(clean_resp)
                print(f"Token count of the response: {token_count}")
                # return generation if generation else "Sorry, I didn't understand that."
                return clean_response(generation) if generation else "Sorry, I didn't understand that."
            except json.JSONDecodeError as e:
                print(f"Error decoding JSON: {str(e)}")
                print("Debug: Response Body String:", response_body_str)  # Debugging line
                return "Sorry, there was an error processing the response."
            
        bedrock_client = initialize_bedrock_client()

            
        

            

        relevant_segment = find_relevant_segment(question, text_segments)
        # prompt = f"Context: {relevant_segment}\n\nQuestion: {question}\n### Answer:"
        # Assuming 'relevant_segment' is the context and 'question' is the actual question being asked.
        response = get_response(file_id,bedrock_client, relevant_segment, question)


        print(f"Chatbot: {response}")

        cluster_id = generate_cluster_id(question, user_identifier, session_start_time)
        print(cluster_id)

        save_to_database(question, response, user_identifier, session_id, session_start_time, cluster_id, user_name, file_id, question_asked_time)

        return JsonResponse({
            'answer': response,
            'source': 'llm',
            'session_id': session_id,
            'question_asked_time': question_asked_time.strftime("%d-%m-%Y %H:%M:%S")
        })

    return HttpResponseBadRequest('Invalid request.')




import hashlib


@api_key_required
def generate_cluster_id(question, user_identifier, session_start_time):
    # Concatenate question, user_identifier, and session_start_time
    concatenated_str = f"{question}{user_identifier}{session_start_time}"
    
    # Hash the concatenated string using SHA-256
    hash_object = hashlib.sha256(concatenated_str.encode())
    
    # Get the hexadecimal representation of the hash
    hex_digest = hash_object.hexdigest()
    
    # Take the first 5 characters of the hexadecimal representation
    shortened_hex_digest = hex_digest[:5]
    
    # Combine the prefix with the shortened hexadecimal representation
    cluster_id = f"CL-ID_{shortened_hex_digest}"
    
    return cluster_id




@api_key_required
def preprocess_text(text):
    lemmatizer = WordNetLemmatizer()
    p = inflect.engine()

    def convert_number_to_words(token):
        if token.isdigit():
            return p.number_to_words(token)
        else:
            return token

    tokens = word_tokenize(text)
    lemmatized_tokens = [lemmatizer.lemmatize(convert_number_to_words(token.lower())) for token in tokens]
    # print('lemmatized_tokens',lemmatized_tokens)
    return ' '.join(lemmatized_tokens)






@api_key_required
def extract_session_number(question):
    # Extracting session number from the question
    match = re.search(r'session\s*(\d+)', question, re.IGNORECASE)
    # print('match',match)
    return int(match.group(1)) if match else None




@api_key_required
def get_answer_from_database(question,user_identifier,session_id,user_name,file_id,question_asked_time):
    try:
        all_entries = ChatSession.objects.filter(user__username=user_name,file_id=file_id)

        
        if not all_entries:
            return None  # No questions in the database, return None
        
        current_question_session_number = extract_session_number(question)
        processed_current_question = preprocess_text(question)
        # print('processed_current_questions',processed_current_question)

        # Filter entries by session number if present
        relevant_entries = [entry for entry in all_entries if extract_session_number(entry.question) == current_question_session_number]
        # print('relevant_entries',relevant_entries)

        if not relevant_entries:
            return None  # No relevant entries found for the specific session

        processed_questions = [preprocess_text(entry.question) for entry in relevant_entries]
        processed_questions.append(processed_current_question)  # Adding the current question for similarity comparison
        # print('processed_questions',processed_questions)

        # Use SentenceTransformer for creating embeddings
        model = SentenceTransformer('paraphrase-MiniLM-L6-v2')
        question_embeddings = model.encode(processed_questions, convert_to_tensor=True)

        # Convert embeddings to NumPy arrays for similarity calculation
        question_embeddings_np = [embedding.cpu().detach().numpy() for embedding in question_embeddings]

        # Calculate cosine similarities
        similarities = cosine_similarity([question_embeddings_np[-1]], question_embeddings_np[:-1])
        # print('similarities',similarities)

        # Find the most similar question
        most_similar_index = similarities.argmax()
        most_similar_answer = relevant_entries[most_similar_index].answer
        verification_count = relevant_entries[most_similar_index].verification_count
        positive_count = relevant_entries[most_similar_index].positive_count
        negative_feedback_count = relevant_entries[most_similar_index].negative_feedback_count

        # Adjust similarity threshold based on context
        similarity_threshold = 0.8 if 'session' in processed_current_question else 0.8

        if similarities[0][most_similar_index] >= similarity_threshold:
            retrieved_entry = relevant_entries[most_similar_index]
            # if not SimilarQuestion.objects.filter(question=retrieved_entry.question, user_session_id=user_identifier).exists():
            retrieved_entry.retrieval_count += 1
            retrieved_entry.save() 
            # existing_similar_question = SimilarQuestion.objects.filter(question=retrieved_entry.question, user_session_id=user_identifier).exists()
            # if not existing_similar_question:
            
            # Create a SimilarQuestion object
            match = SimilarQuestion.objects.create(
                    question=question,
                    answer=retrieved_entry.answer,
                    session_id=session_id,
                    user_session_id=user_identifier,
                    cluster_id=retrieved_entry.cluster_id,
                    session_start_time=retrieved_entry.session_start_time,
                    question_asked_time=question_asked_time,
                    user_name=user_name,
                    file_id=file_id
                )

            return most_similar_answer,verification_count,positive_count,negative_feedback_count,file_id,question_asked_time
        else:
            return None

    except ChatSession.DoesNotExist:
        return None

@api_key_required
def count_positive_feedback(question):
    try:
        # Count the number of positive feedback for the given question
        positive_feedback_count = Feedback.objects.filter(question=question, feedback=True).count()
        return positive_feedback_count
    except Feedback.DoesNotExist:
        return 0
@api_key_required
def count_negative_feedback(question):
    try:
        # Count the number of positive feedback for the given question
        negative_feedback_count = Feedback.objects.filter(question=question, feedback=False).count()
        return negative_feedback_count
    except Feedback.DoesNotExist:
        return 0


import json
from django.db.models import F
from django.db.models import Max
from django.contrib.sessions.models import Session
from django.contrib.sessions.backends.db import SessionStore
@api_key_required
@csrf_exempt
def submit_feedback(request):
    if request.method == "POST":
        try:
            data = json.loads(request.body.decode('utf-8'))
            question = data.get('question')
            feedback = data.get('feedback')

            # Convert feedback to boolean
            feedback_bool = feedback.lower() == "positive"
            
            # Check if the user has already given feedback for this question

            # Create Feedback object
            feedback_obj = Feedback.objects.create(
                user_identifier=request.session.get('user_identifier'),
                session_id=data.get('session_id'),
                question=question,
                answer=data.get('answer'),
                feedback=feedback_bool
            )
            
            # Increase verification count if feedback is positive
            if feedback_bool:
                # Check if session data has processed feedbacks
                processed_feedbacks = request.session.get('processed_feedbacks', [])

                # Check if this feedback has already been processed
                if feedback_obj.id not in processed_feedbacks:
                    # Count positive feedback for this specific question
                    positive_feedback_count = Feedback.objects.filter(question=question, feedback=True, id__gt=feedback_obj.id).count()

                    # Update verification count for the question
                    latest_chat_session = ChatSession.objects.filter(question=question).latest('id')
                    latest_chat_session.positive_count += positive_feedback_count + 1  # Increment by 1 for the current feedback
                    latest_chat_session.save()

                    # Add this feedback to the processed list
                    processed_feedbacks.append(feedback_obj.id)
                    request.session['processed_feedbacks'] = processed_feedbacks

            else:  # If feedback is negative
                # Check if session data has processed negative feedbacks
                processed_negative_feedbacks = request.session.get('processed_negative_feedbacks', [])

                # Check if this negative feedback has already been processed
                if feedback_obj.id not in processed_negative_feedbacks:
                    # Count negative feedback for this specific question
                    negative_feedback_count = Feedback.objects.filter(question=question, feedback=False, id__gt=feedback_obj.id).count()

                    # Update negative feedback count for the question
                    latest_chat_session = ChatSession.objects.filter(question=question).latest('id')
                    latest_chat_session.negative_feedback_count += negative_feedback_count + 1  # Increment by 1 for the current negative feedback
                    latest_chat_session.save()

                    # Add this negative feedback to the processed list
                    processed_negative_feedbacks.append(feedback_obj.id)
                    request.session['processed_negative_feedbacks'] = processed_negative_feedbacks
                else:
                    # Negative feedback already processed, no need to count again
                    pass

            return JsonResponse({'status': 'success'})
        except json.JSONDecodeError:
            return JsonResponse({'status': 'error', 'message': 'Invalid JSON'}, status=400)
    else:
        return JsonResponse({'status': 'error', 'message': 'Method not allowed'}, status=405)





@api_key_required
def save_to_database(question, answer, user_identifier, session_id, session_start_time, cluster_id, user_name,file_id,question_asked_time):
   
    try:
        user = User.objects.get(username=user_name)
    except User.DoesNotExist:
        # Handle the case where the user does not exist
        return None  # or raise an exception

    # Create a ChatSession object
    chat_session = ChatSession.objects.create(
        session_id=session_id,
        session_start_time=session_start_time,
        question=question,
        question_asked_time=question_asked_time,
        answer=answer,
        user_identifier=user_identifier,
        cluster_id=cluster_id,
        user=user,
        file_id=file_id
    )

    # Create a SimilarQuestion object
    similar_question = SimilarQuestion.objects.create(
        session_id=session_id,
        session_start_time=session_start_time,
        question=question,
        question_asked_time=question_asked_time,
        answer=answer,
        user_session_id=user_identifier,
        cluster_id=cluster_id,
        user_name=user_name,
        file_id=file_id
    )

    
    chat_session.save()
    similar_question.save()








from datetime import datetime


from django.http import HttpResponseRedirect
import csv
@api_key_required
def upload_csv(request):
    if request.method == "POST":
        csv_file = request.FILES['csv_file']
        decoded_file = csv_file.read().decode('utf-8').splitlines()
        reader = csv.DictReader(decoded_file)
        for row in reader:
            session_id = row['session_id']
            session_start_time = parse_datetime(row['session_start_time'])
            question = row['question']
            question_asked_time = parse_datetime(row['question_asked_time'])
            answer = row['answer']
            user_identifier = row['user_identifier']
            retrieval_count = int(row['retrieval_count'])
            verification_count = int(row['verification_count'])
            cluster_id = row.get('cluster_id', None)  # Handling if cluster_id is missing in some rows
            
            ChatSession.objects.create(
                session_id=session_id,
                session_start_time=session_start_time,
                question=question,
                question_asked_time=question_asked_time,
                answer=answer,
                user_identifier=user_identifier,
                retrieval_count=retrieval_count,
                verification_count=verification_count,
                cluster_id=cluster_id
            )
        return HttpResponseRedirect('/admin/app1/chatsession/')  # Redirect to admin page after upload
    return render(request, "upload_csv.html")  # Render the upload form if GET request
@api_key_required
def parse_datetime(datetime_str):
    try:
        return datetime.fromisoformat(datetime_str)
    except ValueError:
        # Handle datetime string in alternative format if necessary
        # For example: '2024-04-11 04:24:37+00:00' to '2024-04-11 04:24:37'
        return datetime.strptime(datetime_str.split('+')[0], '%Y-%m-%d %H:%M:%S')





