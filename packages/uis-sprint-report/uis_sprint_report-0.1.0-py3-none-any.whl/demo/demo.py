from rich import print
from rich.table import Table
from pptx import Presentation

from langchain_core.output_parsers import PydanticOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_community.chat_models import ChatOllama
from .prompts import prompts
from .models import SprintReport, ResponseModel


def prepare_parser_and_prompt(model):
    """Prepare parser and prompt template for structured outputs."""
    parser = PydanticOutputParser(pydantic_object=model)
    prompt_template = ChatPromptTemplate.from_messages([
        ("system",
         "Answer the user query. Wrap the output in `json` tags\n{format_instructions}. The context is {context}"),
        ("human", "{query}"),
    ]).partial(format_instructions=parser.get_format_instructions())

    return parser, prompt_template


def prepare_chain(llm, retriever, model):
    parser, prompt_template = prepare_parser_and_prompt(model)
    return {"context": retriever, "query": RunnablePassthrough()} | prompt_template | llm | parser


def invoke_chain(chain, prompt, max_attempts):
    count_attempts = 0
    while count_attempts < max_attempts:
        try:
            response = chain.invoke(prompt)
            return response
        except Exception as e:
            count_attempts += 1
    print("[red]❌ Unable to generate answer[/red]")


def get_report(llm, embeddings, max_attempts, sprint_goals=None):
    chain = prepare_chain(llm, embeddings, SprintReport)
    prompt = prompts['generate_report']['question']
    if sprint_goals is not None:
        prompt += prompts['generate_report']['additional_info'] + sprint_goals
    report = invoke_chain(chain, prompt, max_attempts)
    return report


def report(llm, embeddings, max_attempts):
    print("[bold]Generating report...[/bold]")
    report = get_report(llm, embeddings, max_attempts)

    table = Table(title="Sprint Activities Report")

    table.add_column("Title", style="cyan", no_wrap=True)
    table.add_column("Description", style="magenta")
    table.add_column("Status", style="green")

    for activity in report.activities:
        table.add_row(
            activity.title,
            activity.brief_desc_status,
            activity.status
        )

    print("[green]✅ Report generated[/green]")
    print(table)


def pptx(llm, embeddings, max_attempts, sprint_goals, pptx_file_name):
    report = get_report(llm, embeddings, max_attempts, sprint_goals)

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = "Sprint Activities Report"
    content.text = "Activities Summary:\n"

    for activity in report.activities:
        content.text += f"- {activity.title}: {activity.brief_desc_status} (Status: {activity.status})\n"

    prs.save(pptx_file_name)
    print("[green]✅ PowerPoint report generated[/green]")


def chat(llm, embeddings, max_attempts):
    is_end = False
    while not is_end:
        user_input = input("You: ")
        if user_input == "exit" or user_input == "quit" or user_input == "q":
            is_end = True
            print("[bold]👋 Goodbye![/bold]")
            break
        chain = prepare_chain(llm, embeddings, ResponseModel)
        response = invoke_chain(
            chain,
            user_input + "\n" + prompts['chat']['question'],
            max_attempts
        )
        print(f"Bot: {response}")


def execute_command(command, embeddings, model, max_tokens, max_attempts, sprint_goals, pptx_file_name):
    llm = ChatOllama(model=model, max_tokens=max_tokens)
    if command == "report":
        report(llm, embeddings, max_attempts)
    elif command == "pptx":
        pptx(llm, embeddings, max_attempts, sprint_goals, pptx_file_name)
    elif command == "chat":
        chat(llm, embeddings, max_attempts)
    else:
        print("Invalid command")
        raise Exception("Invalid command")