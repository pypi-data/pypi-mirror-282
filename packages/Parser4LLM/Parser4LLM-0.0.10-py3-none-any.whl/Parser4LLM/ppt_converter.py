from pptx import Presentation
from pptx.presentation import Presentation as PresentationObject
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List, Set
from pathlib import Path
from Parser4LLM.notification import Notification, DefaultNotification

class PPTXConverter:
    def __init__(self, add_images=False, notification: Notification = DefaultNotification()) -> None:
        self.header_handled = False
        self.add_images = add_images
        self.notification = notification
        
    def is_ocr_pptx(self, file_path):
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    # If we can extract text from the shape, it's likely OCR
                    return True
        # If no text could be extracted from any shape, it's likely non-OCR
        return False

    def convert(self, file_path: str) -> str:
        is_ocr = self.is_ocr_pptx(file_path)
        if is_ocr:
            prs = Presentation(file_path)
            md_content = []
            unique_slides: Set[str] = set()

            # Handle header
            if prs.slides and prs.slides[0].placeholders:
                header_content = self._handle_header(prs.slides[0].placeholders)
                if header_content:
                    md_content.append(header_content)

            for i, slide in enumerate(prs.slides):
                slide_md_content: List[str] = []
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:  # type: ignore
                        slide_md_content += self._handle_table(shape.table)
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and self.add_images:  # type: ignore
                        slide_md_content.append(self._handle_image(shape))
                    elif hasattr(shape, "text"):
                        slide_md_content.append(self._handle_paragraph(shape.text))

                slide_md_str = "\n".join(slide_md_content)
                if slide_md_str not in unique_slides:
                    unique_slides.add(slide_md_str)
                    slide_md_str = f"## Slide {i+1}\n{slide_md_str}"
                    md_content.append(slide_md_str)

            return "\n".join(md_content)
        else:
            # call marker-api
            return "No Extracted text"

    def _handle_header(self, placeholders) -> str:
        if not self.header_handled:
            parts = []
            for placeholder in placeholders:
                if placeholder.placeholder_format.idx == 0:  # Title placeholder
                    parts.append(f"# {placeholder.text}")
                elif placeholder.placeholder_format.idx == 1:  # Subtitle placeholder
                    parts.append(f"## {placeholder.text}")
            self.header_handled = True
            return "\n".join(parts)
        return ""

    def _handle_paragraph(self, text: str) -> str:
        # Assuming text is a simple paragraph without complex formatting
        # if text contains letters return text
        if any(c.isalpha() for c in text):
            return text + "\n"
        return ""

    def _handle_image(self, shape) -> str:
        image = shape.image
        image_bytes = image.blob
        image_format = image.ext.lstrip(".")
        image_filename = f"images/image_{shape.shape_id}.{image_format}"
        with open(image_filename, "wb") as f:
            f.write(image_bytes)
        return f"![Image {shape.shape_id}](../{image_filename})"

    def _handle_table(self, table) -> List[str]:
        row_content = []
        for i, row in enumerate(table.rows):
            row_content.append(
                "| " + " | ".join(cell.text.strip() for cell in row.cells) + " |"
            )
            if i == 0:
                row_content.append("|" + "---|" * len(row.cells))
        return row_content

    def save_md(self, md_content: str, file_path: str, encoding: str = "utf-8") -> None:
        with open(file_path, "w", encoding=encoding) as f:
            f.write(md_content)